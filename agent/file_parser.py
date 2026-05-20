"""文件解析引擎

支持解析多种格式文件并提取可读内容，用于对话上下文注入。

支持的格式:
  - 纯文本: .txt .md .csv .json .xml .html .css .js .py .log .yaml .yml
  - PDF: .pdf (PyPDF2)
  - Word: .docx (python-docx)
  - Excel: .xlsx (openpyxl)
  - PowerPoint: .pptx (python-pptx)

策略:
  - 小文件 (< 8KB): 返回完整内容
  - 大文件 (>= 8KB): 返回摘要 + 前 4KB 预览，标注可按需请求原文
"""

import os
import json
import csv as csv_module
import io

SMALL_FILE_THRESHOLD = 8 * 1024
LARGE_PREVIEW_SIZE = 4 * 1024

TEXT_EXTENSIONS = {
    ".txt", ".md", ".csv", ".json", ".xml", ".html", ".css", ".js",
    ".py", ".log", ".yaml", ".yml", ".ini", ".cfg", ".toml", ".sh",
    ".bash", ".zsh", ".sql", ".r", ".rb", ".go", ".rs", ".java",
    ".c", ".cpp", ".h", ".hpp", ".ts", ".tsx", ".jsx", ".vue",
    ".conf", ".env", ".properties",
}


def get_file_type(file_path: str) -> str:
    ext = os.path.splitext(file_path)[1].lower()
    if ext in TEXT_EXTENSIONS:
        return "text"
    if ext == ".pdf":
        return "pdf"
    if ext == ".docx":
        return "docx"
    if ext == ".xlsx":
        return "xlsx"
    if ext == ".pptx":
        return "pptx"
    if ext == ".csv":
        return "csv"
    return "unknown"


def parse_file(file_path: str) -> dict:
    """解析单个文件，返回结构化结果

    Args:
        file_path: 文件绝对路径

    Returns:
        {
            "filename": str,
            "file_type": str,
            "file_size": int,
            "is_large": bool,
            "content": str,         # 全量或摘要 + 预览
            "full_content": str,     # 完整内容（仅小文件）
        }
    """
    filename = os.path.basename(file_path)
    file_type = get_file_type(file_path)

    if not os.path.isfile(file_path):
        return {
            "filename": filename,
            "file_type": file_type,
            "file_size": 0,
            "is_large": False,
            "content": f"[文件不存在: {filename}]",
            "full_content": "",
        }

    file_size = os.path.getsize(file_path)

    try:
        if file_type in ("text", "csv"):
            with open(file_path, "r", encoding="utf-8") as f:
                full = f.read()
        elif file_type == "pdf":
            full = _parse_pdf(file_path)
        elif file_type == "docx":
            full = _parse_docx(file_path)
        elif file_type == "xlsx":
            full = _parse_xlsx(file_path)
        elif file_type == "pptx":
            full = _parse_pptx(file_path)
        else:
            full = f"[不支持的文件格式: .{os.path.splitext(filename)[1]}]"
    except Exception as e:
        return {
            "filename": filename,
            "file_type": file_type,
            "file_size": file_size,
            "is_large": False,
            "content": f"[解析失败: {filename}, 错误: {str(e)}]",
            "full_content": "",
        }

    is_large = len(full.encode("utf-8")) >= SMALL_FILE_THRESHOLD if full else False

    if is_large:
        content = _build_large_file_content(filename, file_type, full, file_size)
    else:
        content = _build_small_file_content(filename, file_type, full)

    return {
        "filename": filename,
        "file_type": file_type,
        "file_size": file_size,
        "is_large": is_large,
        "content": content,
        "full_content": full if not is_large else "",
    }


def parse_files(file_paths: list[str]) -> list[dict]:
    """批量解析文件"""
    return [parse_file(p) for p in file_paths]


def build_context_prompt(parsed_files: list[dict]) -> str:
    """根据解析结果构建注入对话的上下文文本"""
    if not parsed_files:
        return ""

    parts = ["=== 用户上传文件内容 ===\n"]

    for pf in parsed_files:
        parts.append(pf["content"])
        parts.append("")

    parts.append("=== 文件内容结束 ===\n")
    parts.append("请基于以上文件内容回答用户的问题。")
    parts.append("如果用户的问题与文件无关，可以先忽略文件内容，正常回答。")

    return "\n".join(parts)


def generate_file_summary(full_content: str, filename: str) -> str:
    """为文件内容生成一个人类可读的摘要

    用于前端展示上传文件卡片时显示简要信息。
    """
    lines = full_content.strip().split("\n")
    line_count = len(lines)
    char_count = len(full_content)
    size_str = _format_size(char_count)

    first_lines = "\n".join(lines[:3])
    if len(first_lines) > 200:
        first_lines = first_lines[:200] + "..."

    return f"{filename} ({size_str}, {line_count}行)"


def _parse_pdf(file_path: str) -> str:
    try:
        from PyPDF2 import PdfReader
    except ImportError:
        return "[PDF解析需要安装 PyPDF2: pip install PyPDF2]"

    reader = PdfReader(file_path)
    parts = []
    for i, page in enumerate(reader.pages):
        text = page.extract_text()
        if text and text.strip():
            parts.append(f"--- 第{i + 1}页 ---\n{text.strip()}")
    return "\n\n".join(parts) if parts else "[PDF文件中未提取到文本内容，可能是扫描件或图片型PDF]"


def _parse_docx(file_path: str) -> str:
    try:
        from docx import Document
    except ImportError:
        return "[Word解析需要安装 python-docx: pip install python-docx]"

    doc = Document(file_path)
    parts = []
    for para in doc.paragraphs:
        if para.text.strip():
            parts.append(para.text)

    for table in doc.tables:
        table_text = []
        for row in table.rows:
            row_text = " | ".join(cell.text for cell in row.cells)
            table_text.append(row_text)
        parts.append("\n".join(table_text))

    return "\n".join(parts) if parts else "[Word文档内容为空]"


def _parse_xlsx(file_path: str) -> str:
    try:
        from openpyxl import load_workbook
    except ImportError:
        return "[Excel解析需要安装 openpyxl: pip install openpyxl]"

    wb = load_workbook(file_path, data_only=True)
    parts = []

    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        parts.append(f"=== 工作表: {sheet_name} ===")

        rows_list = []
        for row in ws.iter_rows(values_only=True):
            row_values = [str(cell) if cell is not None else "" for cell in row]
            if any(v for v in row_values):
                rows_list.append(" | ".join(row_values))

        if len(rows_list) > 500:
            parts.extend(rows_list[:500])
            parts.append(f"... (共 {len(rows_list)} 行，仅显示前 500 行)")
        else:
            parts.extend(rows_list)

    wb.close()
    return "\n".join(parts) if parts else "[Excel工作簿内容为空]"


def _parse_pptx(file_path: str) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return "[PPT解析需要安装 python-pptx: pip install python-pptx]"

    prs = Presentation(file_path)
    parts = []
    for i, slide in enumerate(prs.slides):
        slide_parts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for paragraph in shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        slide_parts.append(text)
            if shape.has_table:
                table = shape.table
                for row in table.rows:
                    row_text = " | ".join(cell.text for cell in row.cells)
                    slide_parts.append(row_text)
        if slide_parts:
            parts.append(f"--- 幻灯片{i + 1} ---\n" + "\n".join(slide_parts))

    return "\n\n".join(parts) if parts else "[PPT演示文稿内容为空]"


def _build_small_file_content(filename: str, file_type: str, full: str) -> str:
    header = f"[文件: {filename} (类型: {file_type}, 大小: 小文件-完整加载)]"
    return f"{header}\n\n{full}"


def _build_large_file_content(filename: str, file_type: str, full: str, file_size: int) -> str:
    preview = full[:LARGE_PREVIEW_SIZE]
    remaining = len(full) - LARGE_PREVIEW_SIZE
    size_str = _format_size(file_size)

    lines = [
        f"[文件: {filename} (类型: {file_type}, 大小: {size_str})]",
        "",
        "【文件摘要】",
    ]

    content_lines = full.strip().split("\n")
    line_count = len(content_lines)
    char_count = len(full)

    lines.append(f"- 总行数: {line_count}")
    lines.append(f"- 总字符数: {char_count}")
    if file_type == "xlsx":
        lines.append(f"- 包含工作表，已提取文本内容")

    lines.append("")
    lines.append("【前段预览】")
    lines.append(preview)

    if remaining > 0:
        lines.append("")
        lines.append(f"... (剩余 {_format_size(remaining)} 未显示)")
        lines.append("如需查看完整内容，请告知需要哪部分信息。")

    return "\n".join(lines)


def _format_size(size_bytes: int) -> str:
    if size_bytes < 1024:
        return f"{size_bytes}B"
    elif size_bytes < 1024 * 1024:
        return f"{size_bytes / 1024:.1f}KB"
    else:
        return f"{size_bytes / (1024 * 1024):.1f}MB"