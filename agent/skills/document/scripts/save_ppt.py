"""
将内容保存为PPT演示文稿(.pptx)，支持简单模式和排版模式

Args:
    filename: PPT文件名（无需扩展名），如"演示文稿"
    content: 简单模式下的文本内容，用"# 标题"标记每页幻灯片；排版模式下可为空字符串
    formatting: 可选，排版配置JSON对象，控制版式、动画、背景、图表等

Returns:
    str: 生成的文件路径
"""
# DEPENDENCIES: python-pptx
import os
import json


def execute(filename: str, content: str = "", formatting: dict = None) -> str:
    output_dir = os.path.join("document_output", "ppt_output")
    os.makedirs(output_dir, exist_ok=True)

    fmt = None
    if formatting and isinstance(formatting, str):
        try:
            fmt = json.loads(formatting)
        except (json.JSONDecodeError, TypeError):
            pass
    elif formatting and isinstance(formatting, dict):
        fmt = formatting

    if fmt and isinstance(fmt, dict):
        from agent.skills.document.ppt_formatter import create_ppt_presentation
        return create_ppt_presentation(filename, content=content, formatting=fmt, output_dir=output_dir)

    # 简单模式
    from pptx import Presentation
    from pptx.util import Inches, Pt

    filepath = os.path.join(output_dir, f"{filename}.pptx")

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slides_data = []
    current_slide = None
    for line in [l.strip() for l in content.replace("\r", "").split("\n") if l.strip()]:
        if line.startswith("# "):
            if current_slide is not None:
                slides_data.append(current_slide)
            current_slide = {"title": line[2:].strip(), "content": []}
        elif current_slide is not None:
            current_slide["content"].append(line)
    if current_slide is not None:
        slides_data.append(current_slide)

    if not slides_data:
        slides_data = [{"title": filename, "content": [p.strip() for p in content.replace("\r", "").split("\n") if p.strip()]}]

    for slide_data in slides_data:
        slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(slide_layout)
        if slide.shapes.title:
            slide.shapes.title.text = slide_data["title"]
        if slide_data["content"] and len(slide.placeholders) > 1:
            body = slide.placeholders[1]
            tf = body.text_frame
            tf.clear()
            for i, line in enumerate(slide_data["content"]):
                if i == 0:
                    tf.paragraphs[0].text = line
                    tf.paragraphs[0].font.size = Pt(18)
                else:
                    p = tf.add_paragraph()
                    p.text = line
                    p.font.size = Pt(18)

    prs.save(filepath)
    return f"PPT演示文稿已成功保存至: {filepath}"