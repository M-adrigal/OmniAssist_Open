"""
PPT 演示文稿格式化引擎

将结构化的幻灯片描述（JSON）转换为 python-pptx 演示文稿。
支持多页幻灯片、版式、文本样式、形状、图片、表格、图表、背景、动画、过渡。

设计原则：
- 版式可自动推断或用户显式指定
- 所有尺寸单位统一使用厘米(cm)，内部自动转换为英制单位(EMU)
- 未指定的属性使用主题默认值
- 动画使用口语化名称映射到具体 XML 参数
"""

import os
import re
import copy

from pptx import Presentation
from pptx.util import Inches, Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR_TYPE
from pptx.oxml.ns import qn, nsmap
from pptx.oxml import parse_xml
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData

SLIDE_WIDTH_16_9 = Cm(33.867)
SLIDE_HEIGHT_16_9 = Cm(19.05)
SLIDE_WIDTH_4_3 = Cm(25.4)
SLIDE_HEIGHT_4_3 = Cm(19.05)

DEFAULT_THEME = {
    "primary_color": "#4472C4",
    "secondary_color": "#ED7D31",
    "accent_color": "#70AD47",
    "background_color": "#FFFFFF",
    "dark_background": "#1F4E79",
    "title_font": {"name": "微软雅黑", "size": 36, "bold": True, "color": "#1F4E79"},
    "subtitle_font": {"name": "微软雅黑", "size": 20, "color": "#666666"},
    "body_font": {"name": "微软雅黑", "size": 18, "color": "#333333"},
    "footer_font": {"name": "微软雅黑", "size": 10, "color": "#999999"},
    "page_number_font": {"name": "微软雅黑", "size": 10, "color": "#999999"},
}

LAYOUT_PRESETS = {
    "cover": {
        "description": "封面页：居中大标题 + 副标题",
        "elements": [
            {"type": "title", "position": {"left_cm": 3, "top_cm": 5, "width_cm": 27.867, "height_cm": 5}},
            {"type": "subtitle", "position": {"left_cm": 3, "top_cm": 11, "width_cm": 27.867, "height_cm": 3}},
        ],
    },
    "section_header": {
        "description": "章节标题页：居中标题，通常带背景色",
        "elements": [
            {"type": "title", "position": {"left_cm": 3, "top_cm": 6, "width_cm": 27.867, "height_cm": 5}},
        ],
    },
    "toc": {
        "description": "目录页：标题 + 编号列表",
        "elements": [
            {"type": "title", "position": {"left_cm": 3, "top_cm": 2, "width_cm": 27.867, "height_cm": 3}},
            {"type": "text", "position": {"left_cm": 5, "top_cm": 6, "width_cm": 23.867, "height_cm": 11}},
        ],
    },
    "title_content": {
        "description": "标题+正文：顶部标题 + 正文区域",
        "elements": [
            {"type": "title", "position": {"left_cm": 2.5, "top_cm": 1.5, "width_cm": 28.867, "height_cm": 3}},
            {"type": "text", "position": {"left_cm": 3, "top_cm": 5.5, "width_cm": 27.867, "height_cm": 12}},
        ],
    },
    "two_column": {
        "description": "两栏布局：标题 + 左右两栏",
        "elements": [
            {"type": "title", "position": {"left_cm": 2.5, "top_cm": 1.5, "width_cm": 28.867, "height_cm": 3}},
            {"type": "text", "position": {"left_cm": 2, "top_cm": 5.5, "width_cm": 14, "height_cm": 12}},
            {"type": "text", "position": {"left_cm": 17.5, "top_cm": 5.5, "width_cm": 14, "height_cm": 12}},
        ],
    },
    "image_text": {
        "description": "图文混排：标题 + 左图右文",
        "elements": [
            {"type": "title", "position": {"left_cm": 2.5, "top_cm": 1.5, "width_cm": 28.867, "height_cm": 3}},
            {"type": "image", "position": {"left_cm": 2, "top_cm": 5.5, "width_cm": 14, "height_cm": 10}},
            {"type": "text", "position": {"left_cm": 17.5, "top_cm": 5.5, "width_cm": 14, "height_cm": 12}},
        ],
    },
    "blank": {
        "description": "空白页：无预设元素",
        "elements": [],
    },
}

SHAPE_TYPE_MAP = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "circle": MSO_SHAPE.OVAL,
    "oval": MSO_SHAPE.OVAL,
    "arrow": MSO_SHAPE.RIGHT_ARROW,
    "left_arrow": MSO_SHAPE.LEFT_ARROW,
    "up_arrow": MSO_SHAPE.UP_ARROW,
    "down_arrow": MSO_SHAPE.DOWN_ARROW,
    "line": MSO_SHAPE.LINE_INVERSE,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "diamond": MSO_SHAPE.DIAMOND,
    "star": MSO_SHAPE.STAR_5_POINT,
    "矩形": MSO_SHAPE.RECTANGLE,
    "圆角矩形": MSO_SHAPE.ROUNDED_RECTANGLE,
    "圆形": MSO_SHAPE.OVAL,
    "椭圆": MSO_SHAPE.OVAL,
    "箭头": MSO_SHAPE.RIGHT_ARROW,
    "左箭头": MSO_SHAPE.LEFT_ARROW,
    "上箭头": MSO_SHAPE.UP_ARROW,
    "下箭头": MSO_SHAPE.DOWN_ARROW,
    "直线": MSO_SHAPE.LINE_INVERSE,
    "三角形": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "菱形": MSO_SHAPE.DIAMOND,
    "五角星": MSO_SHAPE.STAR_5_POINT,
}

ALIGNMENT_MAP = {
    "left": PP_ALIGN.LEFT,
    "center": PP_ALIGN.CENTER,
    "right": PP_ALIGN.RIGHT,
    "justify": PP_ALIGN.JUSTIFY,
    "左对齐": PP_ALIGN.LEFT,
    "居中": PP_ALIGN.CENTER,
    "右对齐": PP_ALIGN.RIGHT,
    "两端对齐": PP_ALIGN.JUSTIFY,
}

VERTICAL_ANCHOR_MAP = {
    "top": MSO_ANCHOR.TOP,
    "middle": MSO_ANCHOR.MIDDLE,
    "bottom": MSO_ANCHOR.BOTTOM,
    "顶部": MSO_ANCHOR.TOP,
    "中部": MSO_ANCHOR.MIDDLE,
    "底部": MSO_ANCHOR.BOTTOM,
}

ANIMATION_EFFECT_MAP = {
    "entrance": {
        "fly_in": "fly",
        "飞入": "fly",
        "fade": "fade",
        "淡入": "fade",
        "fade_in": "fade",
        "wipe": "wipe",
        "擦除": "wipe",
        "zoom": "zoom",
        "缩放": "zoom",
        "float_in": "float",
        "浮入": "float",
        "split": "split",
        "劈裂": "split",
        "appear": "appear",
        "出现": "appear",
        "random_bars": "randomBars",
        "随机线条": "randomBars",
        "grow": "grow",
        "伸展": "grow",
        "swivel": "swivel",
        "旋转": "swivel",
        "bounce": "bounce",
        "弹跳": "bounce",
    },
    "exit": {
        "fly_out": "fly",
        "飞出": "fly",
        "fade": "fade",
        "淡出": "fade",
        "fade_out": "fade",
        "zoom": "zoom",
        "缩小": "zoom",
        "disappear": "appear",
        "消失": "appear",
    },
    "emphasis": {
        "spin": "spin",
        "旋转": "spin",
        "grow_shrink": "growShrink",
        "放大缩小": "growShrink",
        "pulse": "pulse",
        "脉冲": "pulse",
        "teeter": "teeter",
        "跷跷板": "teeter",
        "color_pulse": "colorPulse",
        "变色": "colorPulse",
        "transparency": "transparency",
        "透明": "transparency",
    },
}

ANIMATION_DIRECTION_MAP = {
    "from_bottom": "fromB",
    "from_top": "fromT",
    "from_left": "fromL",
    "from_right": "fromR",
    "自底部": "fromB",
    "自顶部": "fromT",
    "自左侧": "fromL",
    "自右侧": "fromR",
    "底部": "fromB",
    "顶部": "fromT",
    "左侧": "fromL",
    "右侧": "fromR",
}

TRANSITION_MAP = {
    "fade": "fade",
    "淡入淡出": "fade",
    "push": "push",
    "推进": "push",
    "wipe": "wipe",
    "擦除": "wipe",
    "split": "split",
    "分割": "split",
    "uncover": "uncover",
    "揭开": "uncover",
    "cover": "cover",
    "覆盖": "cover",
    "clock": "clock",
    "时钟": "clock",
    "zoom": "zoom",
    "缩放": "zoom",
    "random": "random",
    "随机": "random",
    "none": None,
    "无": None,
}

CHART_TYPE_MAP = {
    "bar": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "bar3d": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "line": XL_CHART_TYPE.LINE,
    "pie": XL_CHART_TYPE.PIE,
    "柱状图": XL_CHART_TYPE.COLUMN_CLUSTERED,
    "条形图": XL_CHART_TYPE.BAR_CLUSTERED,
    "折线图": XL_CHART_TYPE.LINE,
    "饼图": XL_CHART_TYPE.PIE,
    "饼状图": XL_CHART_TYPE.PIE,
    "面积图": XL_CHART_TYPE.AREA,
}


def _parse_color(hex_color):
    """将 #RRGGBB 颜色转为 RGBColor"""
    if not hex_color or not isinstance(hex_color, str):
        return None
    hex_color = hex_color.lstrip("#")
    if len(hex_color) == 6:
        return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))
    return None


def _cm_to_emu(cm_value):
    """厘米转 EMU"""
    return int(cm_value * 360000)


def _merge_font(base, override):
    """合并字体配置"""
    result = dict(base) if base else {}
    if override:
        result.update({k: v for k, v in override.items() if v is not None})
    return result


def _apply_text_style(text_frame, style, theme, default_font_key="body_font"):
    """对文本框应用文本样式"""
    font_spec = {}
    if style and style.get("font"):
        font_spec = style["font"]

    theme_font = theme.get(default_font_key, {})
    merged_font = _merge_font(theme_font, font_spec)

    alignment = style.get("alignment", "left") if style else "left"
    if alignment in ALIGNMENT_MAP:
        for paragraph in text_frame.paragraphs:
            paragraph.alignment = ALIGNMENT_MAP[alignment]

    vertical = style.get("vertical", "top") if style else "top"
    if vertical in VERTICAL_ANCHOR_MAP:
        text_frame.paragraphs[0].alignment = text_frame.paragraphs[0].alignment

    for paragraph in text_frame.paragraphs:
        _apply_paragraph_font(paragraph, merged_font)
        line_spacing = (style or {}).get("line_spacing")
        if line_spacing:
            paragraph.line_spacing = Pt(line_spacing)


def _apply_paragraph_font(paragraph, font_spec):
    """对段落应用字体"""
    for run in paragraph.runs:
        if "name" in font_spec:
            run.font.name = font_spec["name"]
        if "size" in font_spec:
            run.font.size = Pt(font_spec["size"])
        color = _parse_color(font_spec.get("color"))
        if color:
            run.font.color.rgb = color
        if font_spec.get("bold"):
            run.font.bold = True
        if font_spec.get("italic"):
            run.font.italic = True
        if font_spec.get("underline"):
            run.font.underline = True


def _set_shape_position(shape, position):
    """设置形状位置和大小"""
    if not position:
        return
    left = position.get("left_cm")
    top = position.get("top_cm")
    width = position.get("width_cm")
    height = position.get("height_cm")
    if left is not None:
        shape.left = Cm(left)
    if top is not None:
        shape.top = Cm(top)
    if width is not None:
        shape.width = Cm(width)
    if height is not None:
        shape.height = Cm(height)


def _add_title_element(slide, element, theme, warnings):
    """添加标题元素"""
    text = element.get("text", "")
    position = element.get("position", {})
    style = element.get("style", {})

    left = Cm(position.get("left_cm", 2.5))
    top = Cm(position.get("top_cm", 1.5))
    width = Cm(position.get("width_cm", 28.867))
    height = Cm(position.get("height_cm", 3))

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    _apply_text_style(tf, style, theme, "title_font")

    return txBox


def _add_subtitle_element(slide, element, theme, warnings):
    """添加副标题元素"""
    text = element.get("text", "")
    position = element.get("position", {})
    style = element.get("style", {})

    left = Cm(position.get("left_cm", 3))
    top = Cm(position.get("top_cm", 11))
    width = Cm(position.get("width_cm", 27.867))
    height = Cm(position.get("height_cm", 3))

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = text
    _apply_text_style(tf, style, theme, "subtitle_font")

    return txBox


def _add_text_element(slide, element, theme, warnings):
    """添加正文文本元素，支持多段落和项目符号"""
    text = element.get("text", "")
    position = element.get("position", {})
    style = element.get("style", {})
    bullet = element.get("bullet")

    left = Cm(position.get("left_cm", 3))
    top = Cm(position.get("top_cm", 5.5))
    width = Cm(position.get("width_cm", 27.867))
    height = Cm(position.get("height_cm", 12))

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    lines = text.split("\n") if text else []
    if not lines:
        lines = [""]

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = line

        if bullet:
            level = bullet.get("level", 0)
            p.level = min(level, 8)
            char = bullet.get("character", "\u2022")
            pPr = p._pPr
            if pPr is None:
                pPr = p._p.get_or_add_pPr()
            buChar = pPr.makeelement(qn('a:buChar'), {'char': char})
            pPr.append(buChar)

    _apply_text_style(tf, style, theme, "body_font")

    return txBox


def _add_image_element(slide, element, theme, warnings):
    """添加图片元素"""
    path = element.get("path", "")
    if not path or not os.path.exists(path):
        warnings.append(f"图片文件不存在: {path}")
        return None

    position = element.get("position", {})
    left = Cm(position.get("left_cm", 2))
    top = Cm(position.get("top_cm", 5.5))
    width = position.get("width_cm")
    height = position.get("height_cm")

    try:
        if width and height:
            pic = slide.shapes.add_picture(path, left, top, Cm(width), Cm(height))
        elif width:
            pic = slide.shapes.add_picture(path, left, top, width=Cm(width))
        elif height:
            pic = slide.shapes.add_picture(path, left, top, height=Cm(height))
        else:
            pic = slide.shapes.add_picture(path, left, top, width=Cm(14))

        rotation = element.get("rotation", 0)
        if rotation:
            pic.rotation = rotation

        return pic
    except Exception as e:
        warnings.append(f"插入图片失败: {str(e)}")
        return None


def _add_shape_element(slide, element, theme, warnings):
    """添加形状元素"""
    shape_type = element.get("shape_type", "rectangle")
    if shape_type not in SHAPE_TYPE_MAP:
        warnings.append(f"不支持的形状类型 '{shape_type}'，已降级为矩形")
        shape_type = "rectangle"

    position = element.get("position", {})
    left = Cm(position.get("left_cm", 5))
    top = Cm(position.get("top_cm", 5))
    width = Cm(position.get("width_cm", 5))
    height = Cm(position.get("height_cm", 5))

    shape = slide.shapes.add_shape(SHAPE_TYPE_MAP[shape_type], left, top, width, height)

    fill_spec = element.get("fill")
    if fill_spec:
        color = _parse_color(fill_spec.get("color"))
        if color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = color

    line_spec = element.get("line")
    if line_spec:
        color = _parse_color(line_spec.get("color", "#000000"))
        if color:
            shape.line.color.rgb = color
        width_pt = line_spec.get("width_pt", 1)
        shape.line.width = Pt(width_pt)

    text = element.get("text")
    if text and shape.has_text_frame:
        shape.text_frame.paragraphs[0].text = text
        style = element.get("style", {})
        _apply_text_style(shape.text_frame, style, theme, "body_font")

    rotation = element.get("rotation", 0)
    if rotation:
        shape.rotation = rotation

    return shape


def _add_table_element(slide, element, theme, warnings):
    """添加表格元素"""
    headers = element.get("headers", [])
    rows = element.get("rows", [])
    if not headers and not rows:
        warnings.append("表格缺少数据")
        return None

    num_cols = len(headers) if headers else (len(rows[0]) if rows else 1)
    num_rows = (1 if headers else 0) + len(rows)

    position = element.get("position", {})
    left = Cm(position.get("left_cm", 2))
    top = Cm(position.get("top_cm", 5))
    width = Cm(position.get("width_cm", 29.867))
    height = Cm(position.get("height_cm", 10))

    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
    table = table_shape.table

    style = element.get("style", {})
    header_style = style.get("header", {"bold": True, "bg_color": "#4472C4", "font_color": "#FFFFFF"})

    if headers:
        for i, h in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = str(h)
            for paragraph in cell.text_frame.paragraphs:
                paragraph.alignment = PP_ALIGN.CENTER
                for run in paragraph.runs:
                    run.font.size = Pt(12)
                    run.font.bold = header_style.get("bold", True)
                    h_color = _parse_color(header_style.get("font_color", "#FFFFFF"))
                    if h_color:
                        run.font.color.rgb = h_color
            bg = _parse_color(header_style.get("bg_color", "#4472C4"))
            if bg:
                _set_cell_fill(cell, bg)

    for r_idx, row_data in enumerate(rows):
        row_num = r_idx + (1 if headers else 0)
        for c_idx, value in enumerate(row_data):
            if c_idx < num_cols:
                cell = table.cell(row_num, c_idx)
                cell.text = str(value)
                for paragraph in cell.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.size = Pt(11)

    return table_shape


def _set_cell_fill(cell, color):
    """设置单元格背景色"""
    tcPr = cell._tc.get_or_add_tcPr()
    solidFill = parse_xml(
        f'<a:solidFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:srgbClr val="{str(color)}"/>'
        f'</a:solidFill>'
    )
    tcPr.append(solidFill)


def _add_chart_element(slide, element, theme, warnings):
    """添加图表元素"""
    chart_type = element.get("chart_type", "bar")
    chart_type = CHART_TYPE_MAP.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    data = element.get("data", [])
    if not data or len(data) < 2:
        warnings.append("图表数据不足")
        return None

    chart_data = CategoryChartData()
    headers = data[0]
    chart_data.categories = headers[1:] if len(headers) > 1 else []

    for row in data[1:]:
        series_name = str(row[0]) if row else ""
        values = row[1:] if len(row) > 1 else []
        numeric_values = []
        for v in values:
            try:
                numeric_values.append(float(v))
            except (ValueError, TypeError):
                numeric_values.append(0)
        chart_data.add_series(series_name, numeric_values)

    position = element.get("position", {})
    left = Cm(position.get("left_cm", 2))
    top = Cm(position.get("top_cm", 5))
    width = Cm(position.get("width_cm", 20))
    height = Cm(position.get("height_cm", 12))

    chart_frame = slide.shapes.add_chart(chart_type, left, top, width, height, chart_data)
    chart = chart_frame.chart

    title = element.get("title")
    if title:
        chart.has_title = True
        chart.chart_title.text_frame.paragraphs[0].text = title

    x_label = element.get("x_axis_title")
    if x_label and hasattr(chart, 'category_axis'):
        chart.category_axis.has_title = True
        chart.category_axis.axis_title.text_frame.paragraphs[0].text = x_label

    y_label = element.get("y_axis_title")
    if y_label and hasattr(chart, 'value_axis'):
        chart.value_axis.has_title = True
        chart.value_axis.axis_title.text_frame.paragraphs[0].text = y_label

    legend_pos = element.get("legend_position", "bottom")
    legend_map = {"bottom": 1, "top": 2, "left": 3, "right": 4, "底部": 1, "顶部": 2, "左侧": 3, "右侧": 4}
    if hasattr(chart, 'legend') and chart.legend is not None:
        chart.legend.position = legend_map.get(legend_pos, 1)
        chart.legend.include_in_layout = False

    return chart_frame


def _add_page_number(slide, element, theme, slide_num, warnings):
    """添加页码"""
    position = element.get("position", {})
    style = element.get("style", {})

    left = Cm(position.get("left_cm", 30))
    top = Cm(position.get("top_cm", 17.5))
    width = Cm(position.get("width_cm", 3))
    height = Cm(position.get("height_cm", 1))

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(slide_num)
    p.alignment = PP_ALIGN.RIGHT
    _apply_text_style(tf, style, theme, "page_number_font")

    return txBox


def _add_footer(slide, element, theme, warnings):
    """添加页脚"""
    text = element.get("text", "")
    position = element.get("position", {})
    style = element.get("style", {})

    left = Cm(position.get("left_cm", 2))
    top = Cm(position.get("top_cm", 17.5))
    width = Cm(position.get("width_cm", 20))
    height = Cm(position.get("height_cm", 1))

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = text
    _apply_text_style(tf, style, theme, "footer_font")

    return txBox


ELEMENT_HANDLERS = {
    "title": _add_title_element,
    "subtitle": _add_subtitle_element,
    "text": _add_text_element,
    "image": _add_image_element,
    "shape": _add_shape_element,
    "table": _add_table_element,
    "chart": _add_chart_element,
    "page_number": _add_page_number,
    "footer": _add_footer,
}


def _apply_background(slide, bg_spec, warnings):
    """应用幻灯片背景"""
    if not bg_spec:
        return

    bg_type = bg_spec.get("type", "solid")
    background = slide.background
    fill = background.fill

    if bg_type == "solid":
        color = _parse_color(bg_spec.get("color", "#FFFFFF"))
        if color:
            fill.solid()
            fill.fore_color.rgb = color
    elif bg_type == "gradient":
        try:
            fill.gradient()
            color1 = _parse_color(bg_spec.get("color1", "#FFFFFF"))
            color2 = _parse_color(bg_spec.get("color2", "#4472C4"))
            if color1:
                fill.gradient_stops[0].color.rgb = color1
            if color2:
                fill.gradient_stops[1].color.rgb = color2
            angle = bg_spec.get("angle", 90)
            fill.angle = angle
        except Exception as e:
            warnings.append(f"渐变背景设置失败: {str(e)}")
    elif bg_type == "image":
        path = bg_spec.get("path")
        if path and os.path.exists(path):
            try:
                fill.user_picture(path)
            except Exception as e:
                warnings.append(f"图片背景设置失败: {str(e)}")


def _apply_animation(slide, anim_spec, shape, warnings):
    """应用动画效果"""
    if not shape:
        return

    anim_type = anim_spec.get("type", "entrance")
    effect_name = anim_spec.get("effect", "fade")

    effect_map = ANIMATION_EFFECT_MAP.get(anim_type, {})
    effect = effect_map.get(effect_name)
    if not effect:
        warnings.append(f"不支持的动画效果 '{effect_name}'，已跳过")
        return

    direction = anim_spec.get("direction")
    if direction:
        direction = ANIMATION_DIRECTION_MAP.get(direction, direction)

    trigger = anim_spec.get("trigger", "on_click")
    delay = anim_spec.get("delay_seconds", 0)
    duration = anim_spec.get("duration_seconds", 0.5)

    try:
        shape_id = shape.shape_id
        spid = shape._element.get("spid")

        if anim_type == "entrance":
            _add_entrance_animation(slide, shape_id, spid, effect, direction, trigger, delay, duration)
        elif anim_type == "exit":
            _add_exit_animation(slide, shape_id, spid, effect, direction, trigger, delay, duration)
        elif anim_type == "emphasis":
            _add_emphasis_animation(slide, shape_id, spid, effect, trigger, delay, duration)
    except Exception as e:
        warnings.append(f"动画设置失败: {str(e)}")


def _ensure_timing_root(slide):
    """确保幻灯片有时序根节点"""
    sld = slide._element
    timing = sld.find(qn('p:timing'))
    if timing is None:
        timing = parse_xml(
            f'<p:timing xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'
        )
        sld.append(timing)
    return timing


def _add_entrance_animation(slide, shape_id, spid, effect, direction, trigger, delay, duration):
    """添加入场动画"""
    timing = _ensure_timing_root(slide)

    tnLst = timing.find(qn('p:tnLst'))
    if tnLst is None:
        tnLst = parse_xml(
            f'<p:tnLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'
        )
        timing.append(tnLst)

    delay_str = f"indefinite" if trigger == "on_click" else str(int(delay * 1000))

    anim_xml = f'''<p:animEffect xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
        <p:cBhvr>
            <p:cTn id="{shape_id + 100}" dur="{int(duration * 1000)}" fill="hold">
                <p:stCondLst><p:cond delay="{delay_str}"/></p:stCondLst>
            </p:cTn>
            <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
            <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
        </p:cBhvr>
        <p:anim effect="{effect}"/>
    </p:animEffect>'''

    if direction:
        anim_xml = anim_xml.replace(
            '</p:animEffect>',
            f'<p:animDir>{direction}</p:animDir></p:animEffect>'
        )

    anim_node = parse_xml(anim_xml)
    tnLst.append(anim_node)


def _add_exit_animation(slide, shape_id, spid, effect, direction, trigger, delay, duration):
    """添加退场动画"""
    timing = _ensure_timing_root(slide)

    tnLst = timing.find(qn('p:tnLst'))
    if tnLst is None:
        tnLst = parse_xml(
            f'<p:tnLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'
        )
        timing.append(tnLst)

    delay_str = f"indefinite" if trigger == "on_click" else str(int(delay * 1000))

    anim_xml = f'''<p:animEffect xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" exit="1">
        <p:cBhvr>
            <p:cTn id="{shape_id + 200}" dur="{int(duration * 1000)}" fill="hold">
                <p:stCondLst><p:cond delay="{delay_str}"/></p:stCondLst>
            </p:cTn>
            <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
            <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
        </p:cBhvr>
        <p:anim effect="{effect}"/>
    </p:animEffect>'''

    if direction:
        anim_xml = anim_xml.replace(
            '</p:animEffect>',
            f'<p:animDir>{direction}</p:animDir></p:animEffect>'
        )

    anim_node = parse_xml(anim_xml)
    tnLst.append(anim_node)


def _add_emphasis_animation(slide, shape_id, spid, effect, trigger, delay, duration):
    """添加强调动画"""
    timing = _ensure_timing_root(slide)

    tnLst = timing.find(qn('p:tnLst'))
    if tnLst is None:
        tnLst = parse_xml(
            f'<p:tnLst xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main"/>'
        )
        timing.append(tnLst)

    delay_str = f"indefinite" if trigger == "on_click" else str(int(delay * 1000))

    anim_xml = f'''<p:animEffect xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main">
        <p:cBhvr>
            <p:cTn id="{shape_id + 300}" dur="{int(duration * 1000)}" fill="hold">
                <p:stCondLst><p:cond delay="{delay_str}"/></p:stCondLst>
            </p:cTn>
            <p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>
            <p:attrNameLst><p:attrName>style.effect</p:attrName></p:attrNameLst>
        </p:cBhvr>
        <p:anim effect="{effect}"/>
    </p:animEffect>'''

    anim_node = parse_xml(anim_xml)
    tnLst.append(anim_node)


def _apply_transition(slide, trans_spec, warnings):
    """应用幻灯片过渡效果"""
    if not trans_spec:
        return

    trans_type = trans_spec.get("type", "none")
    trans_type = TRANSITION_MAP.get(trans_type, trans_type)
    if trans_type is None:
        return

    duration = trans_spec.get("duration_seconds", 0.5)

    try:
        sld = slide._element
        trans_elem = parse_xml(
            f'<p:transition xmlns:p="http://schemas.openxmlformats.org/presentationml/2006/main" '
            f'dur="{int(duration * 1000)}">'
            f'<p:{trans_type}/>'
            f'</p:transition>'
        )
        sld.append(trans_elem)
    except Exception as e:
        warnings.append(f"过渡效果设置失败: {str(e)}")


def _infer_layout(elements):
    """根据元素内容推断合适的版式"""
    element_types = [e.get("type") for e in elements]

    has_title = "title" in element_types
    has_subtitle = "subtitle" in element_types
    has_text = "text" in element_types
    has_image = "image" in element_types
    has_table = "table" in element_types
    has_chart = "chart" in element_types

    text_count = element_types.count("text")

    if has_title and has_subtitle and not has_text and not has_image:
        return "cover"
    if has_title and text_count >= 2:
        return "two_column"
    if has_title and has_image and has_text:
        return "image_text"
    if has_title and (has_text or has_table or has_chart):
        return "title_content"
    if has_title and not has_text:
        return "section_header"

    return "title_content"


def _setup_slide(prs, slide_spec, theme, slide_num, warnings):
    """设置单张幻灯片"""
    layout_name = slide_spec.get("layout")
    elements = slide_spec.get("elements", [])

    if not layout_name:
        layout_name = _infer_layout(elements)

    layout_preset = LAYOUT_PRESETS.get(layout_name, LAYOUT_PRESETS["title_content"])

    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    bg_spec = slide_spec.get("background")
    if bg_spec:
        _apply_background(slide, bg_spec, warnings)
    elif layout_name == "section_header":
        _apply_background(slide, {"type": "solid", "color": theme.get("dark_background", "#1F4E79")}, warnings)

    merged_elements = []
    preset_elements = layout_preset.get("elements", [])

    user_element_types = {e.get("type") for e in elements}
    for pe in preset_elements:
        if pe["type"] not in user_element_types:
            merged_elements.append(pe)
    merged_elements.extend(elements)

    shapes = []
    for elem in merged_elements:
        elem_type = elem.get("type")
        handler = ELEMENT_HANDLERS.get(elem_type)
        if handler:
            if elem_type == "page_number":
                shape = handler(slide, elem, theme, slide_num, warnings)
            else:
                shape = handler(slide, elem, theme, warnings)
            shapes.append(shape)
        else:
            warnings.append(f"未知的元素类型 '{elem_type}'，已跳过")

    animations = slide_spec.get("animations", [])
    for anim in animations:
        target_idx = anim.get("target", 0)
        if 0 <= target_idx < len(shapes) and shapes[target_idx] is not None:
            _apply_animation(slide, anim, shapes[target_idx], warnings)

    transition = slide_spec.get("transition")
    if transition:
        _apply_transition(slide, transition, warnings)

    return slide


def create_ppt_presentation(filename, content=None, formatting=None, output_dir=None):
    """创建格式化 PPT 演示文稿

    Args:
        filename: 文件名（不含扩展名）
        content: 简单模式的文本内容（formatting 为 None 时使用，向后兼容）
        formatting: 排版描述 JSON 对象

    Returns:
        str: 结果消息，包含文件路径和可能的警告
    """
    if output_dir is None:
        output_dir = os.path.join("document_output", "ppt_output")
    os.makedirs(output_dir, exist_ok=True)
    filepath = os.path.join(output_dir, f"{filename}.pptx")

    warnings = []
    prs = Presentation()

    slide_size = formatting.get("slide_size", "16:9") if formatting else "16:9"
    if slide_size == "4:3":
        prs.slide_width = SLIDE_WIDTH_4_3
        prs.slide_height = SLIDE_HEIGHT_4_3
    else:
        prs.slide_width = SLIDE_WIDTH_16_9
        prs.slide_height = SLIDE_HEIGHT_16_9

    theme = dict(DEFAULT_THEME)
    if formatting and formatting.get("theme"):
        user_theme = formatting["theme"]
        for key in theme:
            if key in user_theme:
                if isinstance(theme[key], dict):
                    theme[key] = _merge_font(theme[key], user_theme[key])
                else:
                    theme[key] = user_theme[key]

    if formatting and isinstance(formatting, dict):
        slides = formatting.get("slides", [])
        if not slides:
            warnings.append("未指定幻灯片(slides)，生成了空白演示文稿")

        for idx, slide_spec in enumerate(slides):
            _setup_slide(prs, slide_spec, theme, idx + 1, warnings)
    else:
        text = content or ""
        slides_data = []
        current_slide = None
        for line in [l.strip() for l in text.replace("\r", "").split("\n") if l.strip()]:
            if line.startswith("# "):
                if current_slide is not None:
                    slides_data.append(current_slide)
                current_slide = {"title": line[2:].strip(), "content": []}
            elif current_slide is not None:
                current_slide["content"].append(line)
        if current_slide is not None:
            slides_data.append(current_slide)

        if not slides_data:
            slides_data = [{"title": filename, "content": [p.strip() for p in text.replace("\r", "").split("\n") if p.strip()]}]

        for sd in slides_data:
            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            if slide.shapes.title:
                slide.shapes.title.text = sd["title"]
            if sd["content"] and len(slide.placeholders) > 1:
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.clear()
                for i, line in enumerate(sd["content"]):
                    if i == 0:
                        tf.paragraphs[0].text = line
                        tf.paragraphs[0].font.size = Pt(18)
                    else:
                        p = tf.add_paragraph()
                        p.text = line
                        p.font.size = Pt(18)

    prs.save(filepath)

    result = f"PPT演示文稿已成功保存至: {filepath}"
    if warnings:
        result += "\n[提示] " + "; ".join(warnings)
    return result